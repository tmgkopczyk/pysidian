import os
import pptx
import re
from html2text import html2text
from time import sleep


def get_filepaths(directory):
    file_paths = []
    for root, directories, filenames in os.walk(directory):
        if "Lectures" in root:
            for filename in filenames:
                filepath = os.path.join(root, filename)
                file_paths.append(filepath)
    return file_paths


files = get_filepaths(r"D:\Algonquin College\3 - Summer 2024\CST8316 - PC Troubleshooting")


def main():
    extensions = []
    for file in files:
        extension = file.split(".")[-1]
        if extension == "pptx":
            pptx_file = get_pptx_presentation(file)
            if pptx_file is None:
                continue
            else:
                #pass
                #print(pptx_file)
                create_presentation_directory_structure(pptx_file)
        else:
            continue

def get_pptx_slide_content(slide):
    slide_content = {"title": "", "content": [], "pictures": []}
    try:
        _slide_title_start = slide.shapes.title.text
    except AttributeError:
        _slide_title_start = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.strip() == "":
                    continue
                else:
                    _slide_title = shape.text
                    break
    _slide_title = _slide_title_start.strip()
    slide_title = None
    cisco_title_regex = re.match(r"^\S+(?: \S+){0,10}\x0b\S+(?: \S+){0,10}$", _slide_title)
    if cisco_title_regex:
        _title_list = [x for x in _slide_title.split("\x0b") if x.strip() != ""]
        if _title_list[0][-1].islower() and _title_list[1][0].islower():
            slide_title = " ".join([x.strip() for x in _title_list])
            #print(_title_list)
        elif _title_list[0][-1] == ":" and _title_list[1][0].isupper():
            slide_title = " ".join([x.strip() for x in _title_list])
            #print(_title_list)
        else:
            slide_content["title"] = _title_list[-1]
            slide_content["section"] = _title_list[0]
    else:
        pass
    if slide_title is None:
        slide_title = _slide_title
    else:
        slide_title = slide_title
    slide_content["title"] = html2text(slide_title).strip()
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.text == _slide_title_start:
                continue
            else:
                for paragraph in shape.text_frame.paragraphs:
                    if len(paragraph.text.strip()) <= 2 and paragraph.text.isnumeric():
                        continue
                    elif paragraph.text.strip() == "":
                        continue
                    else:
                        if paragraph.level > 0:
                            slide_content["content"].append("\t" * paragraph.level + "- " + paragraph.text.strip())
                        else:
                            slide_content["content"].append(paragraph.text.strip())
        elif shape.has_table:
            # get the table as a list of lists
            table = []
            for row in shape.table.rows:
                table_row = []
                for cell in row.cells:
                    if cell.text.strip() == "":
                        continue
                    else:
                        table_row.append(" ".join(cell.text.split()))
                table.append(table_row)
            # append the table to the slide content as a Markdown table
            table_content = ""
            table_content += "\n"
            for i in range(len(table)):
                # if it is the first row, then add the header and a separator
                if i == 0:
                    table_content += "|"
                    for j in range(len(table[i])):
                        table_content += table[i][j] + "|"
                    table_content += "\n"
                    table_content += "|"
                    for j in range(len(table[i])):
                        table_content += "---|"
                    table_content += "\n"
                else:
                    table_content += "|"
                    for j in range(len(table[i])):
                        table_content += table[i][j] + "|"
                    table_content += "\n"
            slide_content["content"].append(table_content)
        elif hasattr(shape,"image") and hasattr(shape.image, "blob"):
            slide_content["pictures"].append(shape.image.blob)
        else:
            continue
    if slide_content["title"] == "":
        if len(slide_content["content"]) > 0:
            if len(slide_content["content"][0].split("-")) > 1:
                slide_content["title"] = slide_content["content"][0].split("-")[1].strip()
                slide_content["content"][0] = slide_content["content"][0].split("-")[-1].strip()
            slide_content["title"] = slide_content["content"][0]
            slide_content["content"].pop(0)
        else:
            slide_content["title"] = "No title"
    return slide_content


def get_pptx_slides(presentation):
    pptx_slides = []
    for slide in presentation.slides:
        pptx_slide_content = get_pptx_slide_content(slide)
        if pptx_slide_content["title"] == "" or len(pptx_slide_content["content"]) == 0:
            continue
        else:
            pptx_slides.append(pptx_slide_content)
    return pptx_slides
def get_pptx_presentation(file):
    presentation_dict = {"title": "", "slides": []}
    subject = file.split("\\")[3]
    presentation_dict["subject"] = subject
    subject_items = subject.split(" ", 1)
    file_name = file.split("\\")[-1].split(".")[0]
    try:
        prs = pptx.Presentation(file)
        pptx_slides = get_pptx_slides(prs)
        title_slide = pptx_slides[0]
        if title_slide["title"] == subject_items[-1] or title_slide["title"] == subject_items[0]:
            title_slide["title"] = title_slide["content"][0]
            title_slide["content"].pop(0)
        else:
            pass
        presentation_dict["title"] = title_slide["title"]
        presentation_dict["slides"] = pptx_slides[1:]
    except pptx.exc.PackageNotFoundError:
        return None
    return presentation_dict

class ObsidianVault:
    def __init__(self, path):
        self.path = path

    def create_folder(self, folder_path):
        if not os.path.exists(os.path.join(self.path, folder_path)):
            os.mkdir(os.path.join(self.path, folder_path))
        else:
            pass

    def create_file(self, file_path, content):
        try:
            with open(os.path.join(self.path, file_path), "w", encoding="utf-8") as file:
                file.write(content)
        except OSError:
            pass


def create_presentation_directory_structure(presentation):
    invalid_chars = re.compile(r"[\\/:*?<>|]")
    presentation_title = re.sub(r"[\\/:*?<>|]","",presentation['title'])
    presentation_folder = os.path.join(vault.path, presentation_title)
    if not os.path.exists(presentation_folder):
        os.mkdir(presentation_folder)
    else:
        pass
    for slide in presentation['slides']:
        slide_title = re.sub(r"[\\/:*?<>|]","",slide['title'])
        if slide.get("section"):
            slide_section = re.sub(r"[\\/:*?<>|]","", slide["section"])
            section_path = os.path.join(presentation_folder,slide_section)
            if not os.path.exists(section_path):
                os.mkdir(section_path)
            slide_path = os.path.join(section_path,slide_title)
            try:
                with open(f"{slide_path}.md","w",encoding="utf-8") as f:
                    for content in slide["content"]:
                        f.write(content)
                        f.write("\n")
            except OSError:
                pass
            if slide.get("pictures"):
                for p_index, picture in enumerate(slide["pictures"]):
                    try:
                        with open(f"{slide_path}_{p_index}.png","wb") as image:
                            image.write(picture)
                        with open(f"{slide_path}.md","a",encoding="utf-8") as f:
                            f.write("\n")
                            f.write(f"![[{slide_title}_{p_index}.png]]")
                    except OSError:
                        pass
            sleep(1)
        else:
            slide_path = os.path.join(presentation_folder,slide_title)
            try:
                with open(f"{slide_path}.md","w",encoding="utf-8") as f:
                    for content in slide["content"]:
                        f.write(content)
                        f.write("\n")
            except OSError:
                pass
            if slide.get("pictures"):
                for p_index, picture in enumerate(slide["pictures"]):
                    try:
                        with open(f"{slide_path}_{p_index}.png", "wb") as image:
                            image.write(picture)
                        with open(f"{slide_path}.md", "a", encoding="utf-8") as f:
                            f.write("\n")
                            f.write(f"![[{slide_title}_{p_index}.png]]")
                    except OSError:
                        pass
            sleep(1)

vault = ObsidianVault(r"D:\Algonquin College\3 - Summer 2024\CST8316 - PC Troubleshooting\Notes")


if __name__ == "__main__":
    main()